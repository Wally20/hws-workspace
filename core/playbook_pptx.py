from __future__ import annotations

import base64
import os
import re
import unicodedata
from io import BytesIO
from typing import Any, Dict, Iterable, List, Optional, Sequence, Tuple

from PIL import Image, ImageOps, UnidentifiedImageError
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt


SLIDE_WIDTH = 12_192_000
SLIDE_HEIGHT = 6_858_000
DESIGN_WIDTH = 960.0
DESIGN_HEIGHT = 540.0

FONT_REGULAR = "Poppins"
FONT_BOLD = "Poppins"
FONT_BLACK = "Poppins ExtraBold"

WHITE = "#FFFFFF"
INK = "#171717"
MUTED_INK = "#5F5F5F"
DARK = "#111111"
GREEN = "#168736"
FIELD_GREEN = "#159447"
GOLD = "#D6A34F"
SOFT_WHITE = "#F1F1EE"

MAX_IMAGE_BYTES = 25 * 1024 * 1024
MAX_IMAGE_PIXELS = 40_000_000
FIELD_TABLE_PAGE_SIZE = 10


def _emu_x(value: float) -> Emu:
    return Emu(round(float(value) / DESIGN_WIDTH * SLIDE_WIDTH))


def _emu_y(value: float) -> Emu:
    return Emu(round(float(value) / DESIGN_HEIGHT * SLIDE_HEIGHT))


def _emu_w(value: float) -> Emu:
    return _emu_x(value)


def _emu_h(value: float) -> Emu:
    return _emu_y(value)


def _safe_float(value: Any, minimum: float, maximum: float, fallback: float) -> float:
    try:
        number = float(value)
    except (TypeError, ValueError):
        number = fallback
    return max(minimum, min(maximum, number))


def _safe_color(value: Any, fallback: str = WHITE) -> str:
    normalized = str(value or "").strip()
    if re.fullmatch(r"#[0-9a-fA-F]{6}", normalized):
        return normalized.upper()
    return fallback


def _rgb(value: Any, fallback: str = WHITE) -> RGBColor:
    color = _safe_color(value, fallback)[1:]
    return RGBColor(int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))


def _safe_text(value: Any, limit: int = 20_000) -> str:
    text = str(value or "").replace("\r\n", "\n").replace("\r", "\n")
    text = "".join(
        character
        for character in text
        if character in {"\n", "\t"} or ord(character) >= 32
    )
    if len(text) <= limit:
        return text.strip()
    return f"{text[: max(0, limit - 1)].rstrip()}…"


def _plain_key(value: Any) -> str:
    text = unicodedata.normalize("NFKD", _safe_text(value, 1_000))
    text = "".join(character for character in text if not unicodedata.combining(character))
    return re.sub(r"\s+", " ", text.casefold()).strip()


def _chunk(items: Sequence[Any], size: int) -> List[List[Any]]:
    rows = list(items)
    return [rows[index : index + size] for index in range(0, len(rows), size)] or [[]]


def _set_fill_alpha(shape: Any, opacity: float) -> None:
    """Set solid-fill opacity while keeping the resulting shape editable."""
    try:
        solid_fill = shape._element.spPr.solidFill
        if solid_fill is None or not len(solid_fill):
            return
        color_element = solid_fill[0]
        for child in list(color_element):
            if child.tag.endswith("}alpha"):
                color_element.remove(child)
        alpha = OxmlElement("a:alpha")
        alpha.set("val", str(round(max(0.0, min(1.0, opacity)) * 100_000)))
        color_element.append(alpha)
    except (AttributeError, IndexError, TypeError):
        # A solid, opaque fill is still a valid and readable fallback.
        return


def _set_arrow_head(connector: Any) -> None:
    try:
        line_element = connector.line._get_or_add_ln()
        for child in list(line_element):
            if child.tag.endswith("}tailEnd"):
                line_element.remove(child)
        arrow_head = OxmlElement("a:tailEnd")
        arrow_head.set("type", "triangle")
        arrow_head.set("w", "sm")
        arrow_head.set("len", "sm")
        line_element.append(arrow_head)
    except (AttributeError, TypeError):
        return


def _split_long_text(value: Any, maximum: int = 1_000) -> List[str]:
    text = _safe_text(value)
    if not text:
        return [""]
    parts: List[str] = []
    remaining = text
    while len(remaining) > maximum:
        split_at = remaining.rfind(" ", 0, maximum + 1)
        if split_at < maximum // 2:
            split_at = maximum
        parts.append(remaining[:split_at].strip())
        remaining = remaining[split_at:].strip()
    if remaining:
        parts.append(remaining)
    return parts or [""]


def _sort_trainings(trainings: Iterable[Any]) -> List[Dict[str, Any]]:
    indexed = [
        (index, training)
        for index, training in enumerate(trainings)
        if isinstance(training, dict)
    ]

    def sort_key(item: Tuple[int, Dict[str, Any]]) -> Tuple[int, str, int]:
        index, training = item
        training_date = _safe_text(training.get("date"), 40)
        if re.fullmatch(r"\d{4}-\d{2}-\d{2}", training_date):
            return (0, training_date, index)
        return (1, "", index)

    return [training for _index, training in sorted(indexed, key=sort_key)]


def _sort_field_blocks(blocks: Iterable[Any]) -> List[Dict[str, Any]]:
    indexed = [
        (index, block)
        for index, block in enumerate(blocks)
        if isinstance(block, dict) and block.get("type") != "arrow"
    ]

    def sort_key(item: Tuple[int, Dict[str, Any]]) -> Tuple[int, int, str, int]:
        index, block = item
        title_key = _plain_key(block.get("title"))
        order_match = re.match(r"^o\s*0*(\d+)(?:\b|$)", title_key)
        if order_match:
            return (0, int(order_match.group(1)), title_key, index)
        return (1, index, title_key, index)

    return [block for _index, block in sorted(indexed, key=sort_key)]


def _exercise_export_key(block: Dict[str, Any], fallback_index: int) -> str:
    export_key = _safe_text(block.get("sameExerciseKey"), 220)
    if not export_key:
        try:
            exercise_id = int(block.get("exerciseId") or 0)
        except (TypeError, ValueError):
            exercise_id = 0
        exercise_title = _plain_key(block.get("exerciseTitle"))
        export_key = f"exercise:{exercise_id}" if exercise_id else f"title:{exercise_title}"
    if not export_key or export_key == "title:":
        export_key = f"block:{_safe_text(block.get('id'), 80) or fallback_index}"
    return export_key


def _select_exercise_blocks(
    blocks: Iterable[Any],
    *,
    sort_blocks: bool = True,
) -> List[Dict[str, Any]]:
    source = _sort_field_blocks(blocks) if sort_blocks else [
        block for block in blocks if isinstance(block, dict) and block.get("type") != "arrow"
    ]
    exported_keys = set()
    selected: List[Dict[str, Any]] = []
    for index, block in enumerate(source, start=1):
        if not (block.get("exerciseTitle") or block.get("exerciseId")):
            continue
        if block.get("sameExerciseExport"):
            export_key = _exercise_export_key(block, index)
            if export_key in exported_keys:
                continue
            exported_keys.add(export_key)
        selected.append(block)
    return selected


def _field_table_rows(blocks: Iterable[Any]) -> List[Tuple[str, str, str]]:
    rows: List[Tuple[str, str, str]] = []
    listed_exercise_keys = set()
    for index, block in enumerate(_sort_field_blocks(blocks), start=1):
        if block.get("sameExerciseExport") and (block.get("exerciseTitle") or block.get("exerciseId")):
            export_key = _exercise_export_key(block, index)
            if export_key in listed_exercise_keys:
                continue
            listed_exercise_keys.add(export_key)
        rows.append(
            (
                str(index),
                _safe_text(block.get("title"), 500) or f"Blok {index}",
                _safe_text(block.get("exerciseTitle"), 1_000) or "Geen oefening geselecteerd",
            )
        )
    return rows or [("-", "Nog geen blokken", "Nog geen oefening geselecteerd")]


def _decode_image_data_url(value: Any) -> Optional[Tuple[BytesIO, int, int]]:
    data_url = _safe_text(value, MAX_IMAGE_BYTES * 2)
    if not data_url.startswith("data:image/") or "," not in data_url:
        return None
    header, encoded = data_url.split(",", 1)
    if ";base64" not in header.casefold():
        return None
    try:
        raw_bytes = base64.b64decode(encoded, validate=True)
    except (ValueError, base64.binascii.Error):
        return None
    if not raw_bytes or len(raw_bytes) > MAX_IMAGE_BYTES:
        return None

    try:
        with Image.open(BytesIO(raw_bytes)) as source:
            source.load()
            if source.width <= 0 or source.height <= 0:
                return None
            if source.width * source.height > MAX_IMAGE_PIXELS:
                return None
            image = ImageOps.exif_transpose(source).convert("RGBA")
            image.thumbnail((2_000, 1_500), getattr(Image, "Resampling", Image).LANCZOS)
            output = BytesIO()
            image.save(output, format="PNG", optimize=True)
            output.seek(0)
            return output, image.width, image.height
    except (OSError, ValueError, UnidentifiedImageError, Image.DecompressionBombError):
        return None


class _PresentationBuilder:
    def __init__(
        self,
        data: Dict[str, Any],
        background_paths: Sequence[str],
        logo_path: str,
    ) -> None:
        self.data = data if isinstance(data, dict) else {}
        self.background_paths = [
            str(path)
            for path in background_paths
            if path and os.path.isfile(str(path))
        ]
        self.logo_path = str(logo_path or "")
        self.presentation = Presentation()
        self.presentation.slide_width = SLIDE_WIDTH
        self.presentation.slide_height = SLIDE_HEIGHT
        self.blank_layout = self.presentation.slide_layouts[6]
        self.page_index = 0

        title = _safe_text(self.data.get("title"), 500) or "HWS draaiboek"
        self.presentation.core_properties.title = title
        self.presentation.core_properties.subject = "Draaiboek HWS Voetbalschool"
        self.presentation.core_properties.author = "HWS Voetbalschool"
        self.presentation.core_properties.keywords = "HWS, voetbal, draaiboek"

    def _add_shape(
        self,
        slide: Any,
        shape_type: Any,
        x: float,
        y: float,
        width: float,
        height: float,
        *,
        fill: Optional[str] = None,
        opacity: float = 1.0,
        line: Optional[str] = None,
        line_width: float = 1.0,
    ) -> Any:
        shape = slide.shapes.add_shape(
            shape_type,
            _emu_x(x),
            _emu_y(y),
            _emu_w(max(0.1, width)),
            _emu_h(max(0.1, height)),
        )
        if fill is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(fill)
            _set_fill_alpha(shape, opacity)
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = _rgb(line)
            shape.line.width = Pt(max(0.1, line_width))
        return shape

    def _add_text(
        self,
        slide: Any,
        text: Any,
        x: float,
        y: float,
        width: float,
        height: float,
        *,
        size: float = 16,
        color: str = WHITE,
        bold: bool = False,
        black: bool = False,
        align: Any = PP_ALIGN.LEFT,
        valign: Any = MSO_ANCHOR.MIDDLE,
        margin: float = 2.0,
        fit: bool = True,
        limit: int = 20_000,
    ) -> Any:
        shape = slide.shapes.add_textbox(
            _emu_x(x),
            _emu_y(y),
            _emu_w(max(0.1, width)),
            _emu_h(max(0.1, height)),
        )
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = valign
        text_frame.margin_left = _emu_w(margin)
        text_frame.margin_right = _emu_w(margin)
        text_frame.margin_top = _emu_h(margin)
        text_frame.margin_bottom = _emu_h(margin)
        if fit:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        paragraph = text_frame.paragraphs[0]
        paragraph.text = _safe_text(text, limit)
        paragraph.alignment = align
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.0
        if not paragraph.runs:
            paragraph.add_run()
        for run in paragraph.runs:
            run.font.name = FONT_BLACK if black else FONT_BOLD if bold else FONT_REGULAR
            run.font.size = Pt(max(1.0, size))
            run.font.bold = bool(bold or black)
            run.font.color.rgb = _rgb(color)
        return shape

    def _format_shape_text(
        self,
        shape: Any,
        text: Any,
        *,
        size: float,
        color: str = WHITE,
        bold: bool = True,
        align: Any = PP_ALIGN.CENTER,
        margin: float = 2.0,
    ) -> None:
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.margin_left = _emu_w(margin)
        text_frame.margin_right = _emu_w(margin)
        text_frame.margin_top = _emu_h(margin)
        text_frame.margin_bottom = _emu_h(margin)
        paragraph = text_frame.paragraphs[0]
        paragraph.text = _safe_text(text, 2_000)
        paragraph.alignment = align
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        if not paragraph.runs:
            paragraph.add_run()
        for run in paragraph.runs:
            run.font.name = FONT_BOLD if bold else FONT_REGULAR
            run.font.size = Pt(max(1.0, size))
            run.font.bold = bold
            run.font.color.rgb = _rgb(color)

    def _add_panel(
        self,
        slide: Any,
        x: float,
        y: float,
        width: float,
        height: float,
        *,
        fill: str = DARK,
        opacity: float = 0.66,
        line: str = WHITE,
    ) -> Any:
        return self._add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            y,
            width,
            height,
            fill=fill,
            opacity=opacity,
            line=line,
            line_width=0.7,
        )

    def _add_background(self, slide: Any, shade_opacity: float) -> None:
        background_path = (
            self.background_paths[self.page_index % len(self.background_paths)]
            if self.background_paths
            else ""
        )
        if background_path:
            try:
                slide.shapes.add_picture(
                    background_path,
                    Emu(0),
                    Emu(0),
                    width=Emu(SLIDE_WIDTH),
                    height=Emu(SLIDE_HEIGHT),
                )
            except (OSError, ValueError, UnidentifiedImageError):
                self._add_shape(
                    slide,
                    MSO_SHAPE.RECTANGLE,
                    0,
                    0,
                    DESIGN_WIDTH,
                    DESIGN_HEIGHT,
                    fill="#202020",
                )
        else:
            self._add_shape(
                slide,
                MSO_SHAPE.RECTANGLE,
                0,
                0,
                DESIGN_WIDTH,
                DESIGN_HEIGHT,
                fill="#202020",
            )
        if shade_opacity > 0:
            self._add_shape(
                slide,
                MSO_SHAPE.RECTANGLE,
                0,
                0,
                DESIGN_WIDTH,
                DESIGN_HEIGHT,
                fill="#000000",
                opacity=shade_opacity,
            )

    def _new_slide(
        self,
        title: Optional[str] = None,
        *,
        shade_opacity: float = 0.23,
    ) -> Any:
        slide = self.presentation.slides.add_slide(self.blank_layout)
        self._add_background(slide, shade_opacity)
        self.page_index += 1
        if title:
            title_text = _safe_text(title, 1_000).upper()
            title_size = 42 if len(title_text) <= 34 else 34 if len(title_text) <= 58 else 27
            self._add_text(
                slide,
                title_text,
                45,
                22,
                870,
                68,
                size=title_size,
                bold=True,
                black=True,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
                fit=True,
            )
        return slide

    def _add_connector(
        self,
        slide: Any,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        *,
        color: str = WHITE,
        width: float = 1.5,
        arrow: bool = False,
    ) -> Any:
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            _emu_x(x1),
            _emu_y(y1),
            _emu_x(x2),
            _emu_y(y2),
        )
        connector.line.color.rgb = _rgb(color)
        connector.line.width = Pt(max(0.2, width))
        if arrow:
            _set_arrow_head(connector)
        return connector

    def _add_footer(self, slide: Any, training: Optional[Dict[str, Any]]) -> None:
        if not isinstance(training, dict):
            return
        footer = " | ".join(
            part
            for part in (
                _safe_text(training.get("name"), 300) or "Training",
                _safe_text(training.get("dateLabel") or training.get("date"), 300),
            )
            if part
        )
        if footer:
            self._add_text(
                slide,
                footer,
                80,
                514,
                800,
                18,
                size=8,
                align=PP_ALIGN.CENTER,
                color=WHITE,
            )

    def add_cover(self) -> None:
        slide = self._new_slide(shade_opacity=0.30)
        if self.logo_path and os.path.isfile(self.logo_path):
            try:
                slide.shapes.add_picture(
                    self.logo_path,
                    _emu_x(360),
                    _emu_y(35),
                    width=_emu_w(240),
                    height=_emu_h(240),
                )
            except (OSError, ValueError, UnidentifiedImageError):
                pass

        cover_title = _safe_text(self.data.get("coverTitle"), 1_000) or "HWS VOETBALDAG"
        title_size = 40 if len(cover_title) <= 38 else 31 if len(cover_title) <= 65 else 25
        self._add_text(
            slide,
            cover_title.upper(),
            55,
            305,
            850,
            76,
            size=title_size,
            bold=True,
            black=True,
            align=PP_ALIGN.CENTER,
            fit=True,
        )
        self._add_text(
            slide,
            _safe_text(self.data.get("coverMeta"), 2_000).upper(),
            75,
            386,
            810,
            42,
            size=19,
            bold=True,
            align=PP_ALIGN.CENTER,
            fit=True,
        )
        document_title = _safe_text(self.data.get("title"), 1_000)
        if document_title and _plain_key(document_title) != _plain_key(cover_title):
            self._add_text(
                slide,
                document_title,
                105,
                442,
                750,
                38,
                size=12,
                align=PP_ALIGN.CENTER,
                color=SOFT_WHITE,
                fit=True,
            )

    def add_overview(self) -> None:
        slide = self._new_slide("Overzicht")
        is_amateur = self.data.get("playbookType") == "samenwerkende-amateurclubs"
        include_staff = bool(self.data.get("includeStaff", True))
        include_program = bool(self.data.get("includeProgram", True))
        has_contingencies = bool(_safe_text(self.data.get("contingencies")))

        if is_amateur:
            date_sentence = (
                f"voor {_safe_text(self.data.get('cycleCoverLabel'), 120).lower() or 'de cyclus'}: "
                f"{_safe_text(self.data.get('cycleDateRangeLabel'), 500) or 'cyclus nog in te vullen'}"
            )
        else:
            date_sentence = f"op {_safe_text(self.data.get('eventDateLabel'), 300) or 'datum nog in te vullen'}"

        visible_parts = []
        if include_staff:
            visible_parts.append("de taakverdeling")
        if include_program:
            visible_parts.append("het programma")
        visible_parts.append("de veldplattegrond")
        if has_contingencies:
            visible_parts.append("de afspraken voor onvoorziene situaties")
        if len(visible_parts) > 1:
            parts_text = f"{', '.join(visible_parts[:-1])} en {visible_parts[-1]}"
        else:
            parts_text = visible_parts[0]

        intro = (
            f"Dit draaiboek bundelt alle praktische informatie voor de "
            f"{_safe_text(self.data.get('introSubject'), 180) or 'voetbaldag'} bij "
            f"{_safe_text(self.data.get('clubName'), 300) or 'HWS'} {date_sentence}. "
            f"Het document bevat {parts_text}."
        )

        self._add_panel(slide, 145, 118, 670, 346, fill=SOFT_WHITE, opacity=0.86)
        self._add_text(
            slide,
            intro,
            180,
            145,
            600,
            92,
            size=16,
            color=INK,
            valign=MSO_ANCHOR.TOP,
            fit=True,
        )

        if is_amateur:
            details = [
                ("Club", self.data.get("clubName") or "HWS"),
                ("Cyclus", self.data.get("cycleNumber") or "Nog in te vullen"),
                ("Start cyclus", self.data.get("cycleStartDateLabel") or "Nog in te vullen"),
                ("Einde cyclus", self.data.get("cycleEndDateLabel") or "Nog in te vullen"),
            ]
        else:
            details = [
                ("Club", self.data.get("clubName") or "HWS"),
                ("Datum", self.data.get("eventDateLabel") or "Nog in te vullen"),
                ("Locatie", self.data.get("location") or "Nog in te vullen"),
                ("Aanmeldingen", self.data.get("registrationCount") or "0"),
            ]
        for index, (label, value) in enumerate(details):
            column = index % 2
            row = index // 2
            box_x = 180 + column * 315
            box_y = 266 + row * 82
            self._add_panel(slide, box_x, box_y, 285, 62, fill=WHITE, opacity=0.58, line=WHITE)
            self._add_text(
                slide,
                label.upper(),
                box_x + 12,
                box_y + 7,
                261,
                18,
                size=8.5,
                color=INK,
                bold=True,
                valign=MSO_ANCHOR.TOP,
            )
            self._add_text(
                slide,
                value,
                box_x + 12,
                box_y + 25,
                261,
                28,
                size=13,
                color=MUTED_INK,
                valign=MSO_ANCHOR.TOP,
                fit=True,
            )

    def add_training_dates(self, trainings: Sequence[Dict[str, Any]], no_training: Sequence[Dict[str, Any]]) -> None:
        entries: List[Tuple[str, str, str]] = []
        for index, training in enumerate(_sort_trainings(trainings), start=1):
            entries.append(
                (
                    "Training",
                    _safe_text(training.get("dateLabel") or training.get("date"), 500) or "Datum nog in te vullen",
                    _safe_text(training.get("name"), 1_000) or f"Training {index}",
                )
            )
        sorted_no_training = sorted(
            [row for row in no_training if isinstance(row, dict)],
            key=lambda row: _safe_text(row.get("date"), 40),
        )
        for row in sorted_no_training:
            entries.append(
                (
                    "Geen training",
                    _safe_text(row.get("dateLabel") or row.get("date"), 500) or "Datum nog in te vullen",
                    _safe_text(row.get("description"), 1_000) or "Geen training",
                )
            )

        for page_number, page_rows in enumerate(_chunk(entries, 9), start=1):
            title = "Trainingsdata" if page_number == 1 else f"Trainingsdata ({page_number})"
            slide = self._new_slide(title)
            self._add_panel(slide, 65, 112, 830, 377, fill=SOFT_WHITE, opacity=0.86)
            columns = [("Type", 80, 150), ("Datum", 230, 275), ("Omschrijving", 505, 375)]
            for label, x, width in columns:
                self._add_text(
                    slide,
                    label.upper(),
                    x,
                    130,
                    width,
                    22,
                    size=9,
                    color=INK,
                    bold=True,
                )
            rows = page_rows or [("Training", "Datum nog in te vullen", "Nog in te vullen")]
            row_height = 35
            for index, row in enumerate(rows):
                row_y = 160 + index * row_height
                self._add_panel(
                    slide,
                    78,
                    row_y,
                    804,
                    row_height - 4,
                    fill=WHITE,
                    opacity=0.60 if index % 2 == 0 else 0.45,
                    line=WHITE,
                )
                values = ((row[0], 88, 132, True), (row[1], 238, 257, True), (row[2], 513, 357, False))
                for value, x, width, bold in values:
                    self._add_text(
                        slide,
                        value,
                        x,
                        row_y + 2,
                        width,
                        row_height - 8,
                        size=10.5,
                        color=INK if bold else MUTED_INK,
                        bold=bold,
                        fit=True,
                    )

    def add_staff(self, staff: Sequence[Dict[str, Any]]) -> None:
        rows = [row for row in staff if isinstance(row, dict)] or [
            {"name": "Nog in te vullen", "role": "", "setupTask": ""}
        ]
        include_setup = bool(self.data.get("includeStaffSetupTasks", True))
        for page_number, page_rows in enumerate(_chunk(rows, 8), start=1):
            title = "Taakverdeling" if page_number == 1 else f"Taakverdeling ({page_number})"
            slide = self._new_slide(title)
            x = 78
            top = 120
            width = 804
            if include_setup:
                columns = [("Naam", 18, 218), ("Rol", 248, 205), ("Taak bij uitzetten", 465, 321)]
            else:
                columns = [("Naam", 18, 390), ("Rol", 420, 366)]
            self._add_panel(slide, x, top, width, 32, fill=DARK, opacity=0.80, line=DARK)
            for label, offset, column_width in columns:
                self._add_text(
                    slide,
                    label.upper(),
                    x + offset,
                    top + 3,
                    column_width,
                    25,
                    size=9.5,
                    bold=True,
                )
            row_height = 42
            for index, member in enumerate(page_rows):
                row_y = top + 38 + index * row_height
                self._add_panel(
                    slide,
                    x,
                    row_y,
                    width,
                    row_height - 5,
                    fill=DARK,
                    opacity=0.58 if index % 2 == 0 else 0.47,
                    line=WHITE,
                )
                values = [member.get("name") or "-", member.get("role") or "-"]
                if include_setup:
                    values.append(member.get("setupTask") or "-")
                for value, (_label, offset, column_width) in zip(values, columns):
                    self._add_text(
                        slide,
                        value,
                        x + offset,
                        row_y + 2,
                        column_width,
                        row_height - 9,
                        size=12 if offset == 18 else 10.5,
                        bold=offset == 18,
                        color=WHITE if offset == 18 else SOFT_WHITE,
                        fit=True,
                    )

    def _add_program_icon(self, slide: Any, icon_key: Any, center_x: float, center_y: float, size: float) -> None:
        symbols = {
            "clipboard": "✓",
            "flame": "↑",
            "football": "●",
            "utensils": "P",
            "trophy": "★",
            "camera": "○",
            "medical": "+",
            "cones": "△",
            "clock": "◷",
        }
        circle = self._add_shape(
            slide,
            MSO_SHAPE.OVAL,
            center_x - size / 2,
            center_y - size / 2,
            size,
            size,
            fill=WHITE,
            opacity=0.10,
            line=WHITE,
            line_width=1.0,
        )
        self._format_shape_text(
            circle,
            symbols.get(_safe_text(icon_key, 40), "◷"),
            size=max(5, size * 0.37),
            color=WHITE,
            bold=True,
            margin=0,
        )

    def add_program(self, program: Sequence[Dict[str, Any]]) -> None:
        rows = [row for row in program if isinstance(row, dict)] or [
            {"startTime": "", "endTime": "", "activity": "Nog in te vullen", "icon": "clock"}
        ]
        chunk_size = 14 if len(rows) <= 14 else 12
        for page_number, page_rows in enumerate(_chunk(rows, chunk_size), start=1):
            title = "Programma" if page_number == 1 else f"Programma ({page_number})"
            slide = self._new_slide(title)
            row_count = max(1, chunk_size if len(rows) > chunk_size else len(page_rows))
            available_height = 397
            row_height = min(36.5, available_height / row_count)
            row_box_height = max(19, row_height - 3.5)
            top = 108
            for index, item in enumerate(page_rows):
                row_y = top + index * row_height
                self._add_panel(
                    slide,
                    78,
                    row_y,
                    804,
                    row_box_height,
                    fill=DARK,
                    opacity=0.58 if index % 2 == 0 else 0.47,
                    line=WHITE,
                )
                self._add_program_icon(slide, item.get("icon"), 107, row_y + row_box_height / 2, min(25, row_box_height - 5))
                self._add_text(
                    slide,
                    _safe_text(item.get("startTime"), 40) or "--:--",
                    135,
                    row_y + 1,
                    85,
                    row_box_height / 2 + 3,
                    size=11.5,
                    bold=True,
                )
                self._add_text(
                    slide,
                    _safe_text(item.get("endTime"), 40) or "--:--",
                    135,
                    row_y + row_box_height / 2 - 2,
                    85,
                    row_box_height / 2,
                    size=8,
                    color=SOFT_WHITE,
                )
                self._add_text(
                    slide,
                    item.get("activity") or "Nog in te vullen",
                    235,
                    row_y + 2,
                    625,
                    row_box_height - 4,
                    size=13 if row_height >= 31 else 10,
                    bold=True,
                    fit=True,
                )

    def add_contingencies(self, contingencies: Any) -> None:
        entries: List[Tuple[str, str]] = []
        for line in _safe_text(contingencies).splitlines():
            stripped = line.strip()
            if not stripped:
                continue
            if ":" in stripped:
                label, value = stripped.split(":", 1)
            else:
                label, value = "Scenario", stripped
            pieces = _split_long_text(value, 1_000)
            for piece_index, piece in enumerate(pieces):
                continued_label = label if piece_index == 0 else f"{label} (vervolg)"
                entries.append((_safe_text(continued_label, 500), piece))
        if not entries:
            return

        for page_number, page_rows in enumerate(_chunk(entries, 7), start=1):
            title = "Onvoorziene omstandigheden" if page_number == 1 else f"Onvoorziene omstandigheden ({page_number})"
            slide = self._new_slide(title)
            self._add_panel(slide, 118, 115, 724, 370, fill=SOFT_WHITE, opacity=0.88)
            row_height = 330 / max(1, len(page_rows))
            for index, (label, value) in enumerate(page_rows):
                row_y = 137 + index * row_height
                show_label = bool(label and _plain_key(label) not in {"scenario", "algemeen"})
                if show_label:
                    self._add_text(
                        slide,
                        label.upper(),
                        153,
                        row_y,
                        654,
                        17,
                        size=8.5,
                        color=MUTED_INK,
                        bold=True,
                        valign=MSO_ANCHOR.TOP,
                    )
                self._add_text(
                    slide,
                    value or "Nog in te vullen",
                    153,
                    row_y + (17 if show_label else 0),
                    654,
                    max(18, row_height - (20 if show_label else 4)),
                    size=13,
                    color=INK,
                    valign=MSO_ANCHOR.TOP,
                    fit=True,
                )

    def add_training_cover(self, training: Dict[str, Any], fallback_index: int) -> None:
        slide = self._new_slide(shade_opacity=0.30)
        training_name = _safe_text(training.get("name"), 1_000) or f"Training {fallback_index}"
        date_label = _safe_text(training.get("dateLabel") or training.get("date"), 500) or "Datum nog in te vullen"
        self._add_text(
            slide,
            training_name.upper(),
            80,
            210,
            800,
            90,
            size=44 if len(training_name) <= 36 else 31,
            bold=True,
            black=True,
            align=PP_ALIGN.CENTER,
            fit=True,
        )
        self._add_text(
            slide,
            date_label.upper(),
            100,
            310,
            760,
            44,
            size=20,
            bold=True,
            align=PP_ALIGN.CENTER,
            fit=True,
        )
        age_groups = training.get("ageGroups") if isinstance(training.get("ageGroups"), list) else []
        age_label = " · ".join(_safe_text(age, 40) for age in age_groups if _safe_text(age, 40))
        if age_label:
            self._add_text(
                slide,
                age_label,
                160,
                365,
                640,
                30,
                size=12,
                align=PP_ALIGN.CENTER,
                color=SOFT_WHITE,
            )

    def _add_pitch(self, slide: Any, blocks: Sequence[Dict[str, Any]], x: float, y: float, width: float, height: float) -> None:
        self._add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x - 7,
            y - 7,
            width + 14,
            height + 14,
            fill=GREEN,
            line=None,
        )
        self._add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            x,
            y,
            width,
            height,
            fill=GREEN,
            line=WHITE,
            line_width=1.6,
        )
        self._add_connector(slide, x, y + height / 2, x + width, y + height / 2, width=1.4)
        self._add_shape(
            slide,
            MSO_SHAPE.OVAL,
            x + width * 0.32,
            y + height * 0.38,
            width * 0.36,
            height * 0.24,
            fill=None,
            line=WHITE,
            line_width=1.4,
        )
        self._add_shape(slide, MSO_SHAPE.OVAL, x + width / 2 - 2, y + height / 2 - 2, 4, 4, fill=WHITE)

        penalty_width = width * 0.52
        penalty_height = height * 0.14
        goal_box_width = width * 0.24
        goal_box_height = height * 0.055
        for top in (y, y + height - penalty_height):
            self._add_shape(
                slide,
                MSO_SHAPE.RECTANGLE,
                x + (width - penalty_width) / 2,
                top,
                penalty_width,
                penalty_height,
                fill=None,
                line=WHITE,
                line_width=1.1,
            )
        for top in (y, y + height - goal_box_height):
            self._add_shape(
                slide,
                MSO_SHAPE.RECTANGLE,
                x + (width - goal_box_width) / 2,
                top,
                goal_box_width,
                goal_box_height,
                fill=None,
                line=WHITE,
                line_width=1.0,
            )

        for arrow in [item for item in blocks if isinstance(item, dict) and item.get("type") == "arrow"]:
            self._add_connector(
                slide,
                x + width * (_safe_float(arrow.get("x1"), 0, 100, 20) / 100),
                y + height * (_safe_float(arrow.get("y1"), 0, 100, 25) / 100),
                x + width * (_safe_float(arrow.get("x2"), 0, 100, 60) / 100),
                y + height * (_safe_float(arrow.get("y2"), 0, 100, 35) / 100),
                color=_safe_color(arrow.get("color"), WHITE),
                width=max(0.7, min(4.0, _safe_float(arrow.get("strokeWidth"), 2, 12, 5) * 0.35)),
                arrow=True,
            )

        block_rows = [item for item in blocks if isinstance(item, dict) and item.get("type") != "arrow"]
        for index, block in enumerate(block_rows, start=1):
            block_width = width * (_safe_float(block.get("width"), 2, 100, 20) / 100)
            block_height = height * (_safe_float(block.get("height"), 2, 100, 14) / 100)
            block_x = x + width * (_safe_float(block.get("x"), 0, 100, 8) / 100)
            block_y = y + height * (_safe_float(block.get("y"), 0, 100, 8) / 100)
            if block_x + block_width > x + width:
                block_width = max(3, x + width - block_x)
            if block_y + block_height > y + height:
                block_height = max(3, y + height - block_y)
            shape = self._add_shape(
                slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                block_x,
                block_y,
                block_width,
                block_height,
                fill=_safe_color(block.get("color"), "#D5EFD3"),
                opacity=0.72,
                line=WHITE,
                line_width=0.8,
            )
            block_title = _safe_text(block.get("title"), 500) or f"Blok {index}"
            exercise_title = _safe_text(block.get("exerciseTitle"), 500)
            label = block_title.upper()
            if exercise_title:
                label = f"{label}\n{exercise_title}"
            self._format_shape_text(
                shape,
                label,
                size=max(3.5, min(8.5, block_width * 0.10, block_height * 0.23)),
                color=WHITE,
                bold=True,
                margin=1.5,
            )

    def add_field_layout(
        self,
        blocks: Sequence[Dict[str, Any]],
        *,
        training: Optional[Dict[str, Any]] = None,
        period: Optional[Dict[str, Any]] = None,
    ) -> None:
        normalized_blocks = [block for block in blocks if isinstance(block, dict)]
        table_rows = _field_table_rows(normalized_blocks)
        row_pages = _chunk(table_rows, FIELD_TABLE_PAGE_SIZE)
        for page_number, page_rows in enumerate(row_pages, start=1):
            title = "Veldplattegrond" if page_number == 1 else f"Veldplattegrond ({page_number})"
            slide = self._new_slide(title)
            period_label = ""
            if isinstance(period, dict):
                label = _safe_text(period.get("label"), 500)
                start_time = _safe_text(period.get("startTime"), 50)
                end_time = _safe_text(period.get("endTime"), 50)
                time_label = f"{start_time or '--:--'}-{end_time or '--:--'}" if start_time or end_time else ""
                period_label = " | ".join(part for part in (label, time_label) if part)
            if period_label:
                self._add_text(slide, period_label, 65, 91, 830, 25, size=11, bold=True, align=PP_ALIGN.CENTER)

            self._add_pitch(slide, normalized_blocks, 70, 128, 250, 350)
            table_x = 350
            table_y = 128
            table_width = 550
            columns = [("#", 0, 42), ("Naam blok", 42, 155), ("Oefening", 197, 353)]
            self._add_panel(slide, table_x, table_y, table_width, 29, fill=DARK, opacity=0.82, line=DARK)
            for label, offset, width in columns:
                self._add_text(
                    slide,
                    label.upper(),
                    table_x + offset + 7,
                    table_y + 3,
                    width - 12,
                    22,
                    size=8.5,
                    bold=True,
                )
            row_height = 33
            for index, row in enumerate(page_rows):
                row_y = table_y + 34 + index * row_height
                self._add_panel(
                    slide,
                    table_x,
                    row_y,
                    table_width,
                    row_height - 4,
                    fill=DARK,
                    opacity=0.58 if index % 2 == 0 else 0.46,
                    line=WHITE,
                )
                values = ((row[0], 0, 42, True), (row[1], 42, 155, True), (row[2], 197, 353, False))
                for value, offset, width, bold in values:
                    self._add_text(
                        slide,
                        value,
                        table_x + offset + 7,
                        row_y + 2,
                        width - 12,
                        row_height - 8,
                        size=9.5 if bold else 9,
                        bold=bold,
                        color=WHITE if bold else SOFT_WHITE,
                        fit=True,
                    )
            self._add_footer(slide, training)

    def _add_exercise_field_elements(
        self,
        slide: Any,
        field: Dict[str, Any],
        x: float,
        y: float,
        width: float,
        height: float,
    ) -> None:
        raw_viewbox = field.get("viewBox")
        if not isinstance(raw_viewbox, list) or len(raw_viewbox) != 4:
            raw_viewbox = [0, 0, 100, 70]
        viewbox = [_safe_float(value, -1e12, 1e12, 0) for value in raw_viewbox]
        if viewbox[2] <= 0 or viewbox[3] <= 0:
            viewbox = [0, 0, 100, 70]

        raw_elements = field.get("elements") if isinstance(field.get("elements"), list) else []
        scale = min(width / viewbox[2], height / viewbox[3])
        draw_width = viewbox[2] * scale
        draw_height = viewbox[3] * scale
        offset_x = x + (width - draw_width) / 2 - viewbox[0] * scale
        offset_y = y + (height - draw_height) / 2 - viewbox[1] * scale

        for element in raw_elements[:140]:
            if not isinstance(element, dict):
                continue
            element_x = _safe_float(element.get("x"), -1e12, 1e12, 0)
            element_y = _safe_float(element.get("y"), -1e12, 1e12, 0)
            element_width = max(1, _safe_float(element.get("width"), 0, 1e12, 1))
            element_height = max(1, _safe_float(element.get("height"), 0, 1e12, 1))
            mapped_x = offset_x + element_x * scale
            mapped_y = offset_y + element_y * scale
            mapped_width = max(1.5, element_width * scale)
            mapped_height = max(1.5, element_height * scale)
            fill = _safe_color(element.get("fill"), DARK)
            element_type = _safe_text(element.get("type"), 40)
            if element_type == "line":
                self._add_connector(
                    slide,
                    mapped_x,
                    mapped_y,
                    mapped_x + mapped_width,
                    mapped_y + mapped_height,
                    color=fill,
                    width=1.6,
                )
                continue
            shape_type = MSO_SHAPE.OVAL if element_type == "ellipse" else MSO_SHAPE.TRAPEZOID if element_type == "cone" else MSO_SHAPE.RECTANGLE
            outline = WHITE if fill in {"#000000", "#00B050"} else DARK
            self._add_shape(
                slide,
                shape_type,
                mapped_x,
                mapped_y,
                mapped_width,
                mapped_height,
                fill=fill,
                line=outline,
                line_width=0.7,
            )

        self._add_exercise_overlays(slide, field.get("overlayItems"), x, y, width, height)

    def _add_exercise_overlays(
        self,
        slide: Any,
        raw_items: Any,
        x: float,
        y: float,
        width: float,
        height: float,
    ) -> None:
        items = raw_items if isinstance(raw_items, list) else []
        for item in items[:160]:
            if not isinstance(item, dict):
                continue
            item_type = _safe_text(item.get("type"), 40)
            color = _safe_color(item.get("color"), DARK)
            center_x = x + width * (_safe_float(item.get("x"), 0, 100, 50) / 100)
            center_y = y + height * (_safe_float(item.get("y"), 0, 100, 50) / 100)
            item_scale = _safe_float(item.get("size"), 45, 220, 100) / 100
            if item_type == "player":
                self._add_shape(
                    slide,
                    MSO_SHAPE.OVAL,
                    center_x - 5 * item_scale,
                    center_y - 7 * item_scale,
                    10 * item_scale,
                    10 * item_scale,
                    fill=color,
                    line=WHITE,
                    line_width=0.6,
                )
                self._add_shape(
                    slide,
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    center_x - 5 * item_scale,
                    center_y + 3 * item_scale,
                    10 * item_scale,
                    6 * item_scale,
                    fill=color,
                    line=None,
                )
            elif item_type == "ball":
                self._add_shape(
                    slide,
                    MSO_SHAPE.OVAL,
                    center_x - 4 * item_scale,
                    center_y - 4 * item_scale,
                    8 * item_scale,
                    8 * item_scale,
                    fill=color,
                    line=DARK,
                    line_width=0.7,
                )
            elif item_type in {"cone", "small-cone", "big-cone"}:
                cone_scale = item_scale * (1.3 if item_type == "big-cone" else 1)
                self._add_shape(
                    slide,
                    MSO_SHAPE.ISOSCELES_TRIANGLE,
                    center_x - 5 * cone_scale,
                    center_y - 6 * cone_scale,
                    10 * cone_scale,
                    12 * cone_scale,
                    fill=color,
                    line=DARK,
                    line_width=0.6,
                )
            elif item_type == "goal":
                goal = self._add_shape(
                    slide,
                    MSO_SHAPE.RECTANGLE,
                    center_x - 14 * item_scale,
                    center_y - 8 * item_scale,
                    28 * item_scale,
                    16 * item_scale,
                    fill=None,
                    line=color,
                    line_width=max(0.8, 1.3 * item_scale),
                )
                goal.fill.background()
            elif item_type in {"line", "arrow"}:
                end_x = x + width * (_safe_float(item.get("x2"), 0, 100, 62) / 100)
                end_y = y + height * (_safe_float(item.get("y2"), 0, 100, 50) / 100)
                self._add_connector(
                    slide,
                    center_x,
                    center_y,
                    end_x,
                    end_y,
                    color=color,
                    width=max(0.7, 1.4 * item_scale),
                    arrow=item_type == "arrow",
                )
            elif item_type == "text":
                label = _safe_text(item.get("text"), 200) or "Tekst"
                label_width = min(95, max(35, 8 + len(label) * 4.2)) * item_scale
                panel = self._add_shape(
                    slide,
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    center_x - label_width / 2,
                    center_y - 8 * item_scale,
                    label_width,
                    16 * item_scale,
                    fill=WHITE,
                    opacity=0.90,
                    line=None,
                )
                self._format_shape_text(panel, label, size=max(4, 6.5 * item_scale), color=color, margin=1)

    def _add_exercise_field(
        self,
        slide: Any,
        field_value: Any,
        x: float,
        y: float,
        width: float,
        height: float,
    ) -> None:
        self._add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            y,
            width,
            height,
            fill=FIELD_GREEN,
            line=WHITE,
            line_width=0.8,
        )
        field = field_value if isinstance(field_value, dict) else {}
        decoded = _decode_image_data_url(field.get("imageDataUrl"))
        stage = (x + 7, y + 7, width - 14, height - 14)
        if decoded is not None:
            image_stream, image_width, image_height = decoded
            scale = min(stage[2] / image_width, stage[3] / image_height)
            picture_width = image_width * scale
            picture_height = image_height * scale
            picture_x = stage[0] + (stage[2] - picture_width) / 2
            picture_y = stage[1] + (stage[3] - picture_height) / 2
            try:
                slide.shapes.add_picture(
                    image_stream,
                    _emu_x(picture_x),
                    _emu_y(picture_y),
                    width=_emu_w(picture_width),
                    height=_emu_h(picture_height),
                )
                self._add_exercise_overlays(
                    slide,
                    field.get("overlayItems"),
                    picture_x,
                    picture_y,
                    picture_width,
                    picture_height,
                )
                return
            except (OSError, ValueError, UnidentifiedImageError):
                pass
        self._add_exercise_field_elements(slide, field, *stage)

    def _add_detail_panel(
        self,
        slide: Any,
        label: str,
        value: Any,
        x: float,
        y: float,
        width: float,
        height: float,
    ) -> None:
        self._add_panel(slide, x, y, width, height, fill=DARK, opacity=0.56, line=WHITE)
        self._add_text(
            slide,
            label.upper(),
            x + 10,
            y + 7,
            width - 20,
            18,
            size=7.5,
            color=SOFT_WHITE,
            bold=True,
            valign=MSO_ANCHOR.TOP,
        )
        self._add_text(
            slide,
            _safe_text(value) or "Niet ingevuld",
            x + 10,
            y + 27,
            width - 20,
            max(12, height - 35),
            size=9.5,
            color=WHITE,
            valign=MSO_ANCHOR.TOP,
            fit=True,
        )

    def add_exercise_detail(
        self,
        block: Dict[str, Any],
        fallback_index: int,
        training: Optional[Dict[str, Any]] = None,
    ) -> None:
        block_title = _safe_text(block.get("title"), 500) or f"Blok {fallback_index}"
        exercise_title = _safe_text(block.get("exerciseTitle"), 1_000) or "Oefening"
        slide = self._new_slide(f"{block_title} - {exercise_title}")
        details = block.get("exerciseDetails") if isinstance(block.get("exerciseDetails"), dict) else {}
        exercise_kind = _safe_text(block.get("exerciseKind") or block.get("category"), 500)

        left_x = 45
        center_x = 345
        right_x = 645
        column_width = 270
        top_y = 122
        top_height = 236
        bottom_y = 372
        bottom_height = 122
        variation_height = 112

        self._add_detail_panel(
            slide,
            "Omschrijving oefening",
            details.get("description"),
            left_x,
            top_y,
            column_width,
            top_height,
        )
        self._add_exercise_field(
            slide,
            block.get("exerciseField"),
            center_x,
            top_y,
            column_width,
            top_height,
        )
        self._add_detail_panel(
            slide,
            "Variatie makkelijker maken",
            details.get("variationEasier"),
            right_x,
            top_y,
            column_width,
            variation_height,
        )
        self._add_detail_panel(
            slide,
            "Variatie moeilijker maken",
            details.get("variationHarder"),
            right_x,
            top_y + variation_height + 12,
            column_width,
            variation_height,
        )
        coaching_label = f"Coaching - {exercise_kind}" if exercise_kind else "Coaching"
        self._add_detail_panel(
            slide,
            coaching_label,
            details.get("coaching"),
            left_x,
            bottom_y,
            column_width,
            bottom_height,
        )
        self._add_detail_panel(
            slide,
            "Materialen",
            details.get("materials"),
            center_x,
            bottom_y,
            column_width,
            bottom_height,
        )
        self._add_detail_panel(
            slide,
            "Afmetingen",
            details.get("dimensions"),
            right_x,
            bottom_y,
            column_width,
            bottom_height,
        )
        self._add_footer(slide, training)

    def build(self) -> bytes:
        self.add_cover()
        self.add_overview()

        field_trainings = [
            training
            for training in (self.data.get("fieldTrainings") if isinstance(self.data.get("fieldTrainings"), list) else [])
            if isinstance(training, dict)
        ]
        no_training = [
            row
            for row in (self.data.get("cycleNoTrainingDates") if isinstance(self.data.get("cycleNoTrainingDates"), list) else [])
            if isinstance(row, dict)
        ]
        is_amateur = self.data.get("playbookType") == "samenwerkende-amateurclubs"
        if is_amateur and (len(field_trainings) > 1 or no_training):
            self.add_training_dates(field_trainings, no_training)

        if bool(self.data.get("includeStaff", True)):
            staff = self.data.get("staff") if isinstance(self.data.get("staff"), list) else []
            self.add_staff(staff)
        if bool(self.data.get("includeProgram", True)):
            program = self.data.get("program") if isinstance(self.data.get("program"), list) else []
            self.add_program(program)
        contingencies = _safe_text(self.data.get("contingencies"))
        if contingencies:
            self.add_contingencies(contingencies)

        if is_amateur and field_trainings:
            for training_index, training in enumerate(field_trainings, start=1):
                self.add_training_cover(training, training_index)
                periods = training.get("fieldPeriods") if isinstance(training.get("fieldPeriods"), list) else []
                periods = [period for period in periods if isinstance(period, dict)]
                if not periods:
                    periods = [
                        {
                            "label": "Plattegrond 1",
                            "startTime": "",
                            "endTime": "",
                            "fieldLayout": training.get("fieldLayout") if isinstance(training.get("fieldLayout"), list) else [],
                        }
                    ]
                all_training_blocks: List[Dict[str, Any]] = []
                for period in periods:
                    blocks = period.get("fieldLayout") if isinstance(period.get("fieldLayout"), list) else []
                    normalized_blocks = [block for block in blocks if isinstance(block, dict)]
                    self.add_field_layout(normalized_blocks, training=training, period=period)
                    all_training_blocks.extend(_sort_field_blocks(normalized_blocks))
                for exercise_index, block in enumerate(
                    _select_exercise_blocks(all_training_blocks, sort_blocks=False),
                    start=1,
                ):
                    self.add_exercise_detail(block, exercise_index, training)
        else:
            blocks = self.data.get("fieldLayout") if isinstance(self.data.get("fieldLayout"), list) else []
            normalized_blocks = [block for block in blocks if isinstance(block, dict)]
            self.add_field_layout(normalized_blocks)
            for exercise_index, block in enumerate(_select_exercise_blocks(normalized_blocks), start=1):
                self.add_exercise_detail(block, exercise_index)

        buffer = BytesIO()
        self.presentation.save(buffer)
        buffer.seek(0)
        content = buffer.read()
        if not content.startswith(b"PK"):
            raise RuntimeError("De PowerPoint-export leverde geen geldig presentatiebestand op.")
        return content


def create_playbook_presentation(
    data: Dict[str, Any],
    background_paths: Sequence[str],
    logo_path: str,
) -> bytes:
    """Build a native, editable 16:9 PowerPoint for a normalized HWS playbook."""
    try:
        builder = _PresentationBuilder(
            data if isinstance(data, dict) else {},
            background_paths if isinstance(background_paths, (list, tuple)) else [],
            logo_path,
        )
        return builder.build()
    except RuntimeError:
        raise
    except Exception as exc:
        raise RuntimeError(
            "De PowerPoint-export kon niet worden gemaakt. Controleer de draaiboekinhoud "
            "en eventuele oefeningafbeeldingen en probeer het opnieuw."
        ) from exc
